from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Understanding von Neumann and Wittgenstein"
subtitle.text = "A Philosophical and Computational Perspective"

# Slide 2: von Neumann's Quote
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]
title.text = "John von Neumann's Perspective"
content.text = (
    "Quote:\n"
    "\"You insist that there is something a machine cannot do. If you will tell me precisely what it is that a machine cannot do, then I can always make a machine which will do just that!\"\n"
    "- John von Neumann\n\n"
    "Key Points:\n"
    "1. Expresses confidence in the capabilities of machines (computers).\n"
    "2. Reflects a belief in the continuous advancement of technology.\n"
    "3. Highlights the theoretical vs. practical aspects of computational problems."
)

# Slide 3: Wittgenstein's Quote
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]
title.text = "Ludwig Wittgenstein's Perspective"
content.text = (
    "Quote:\n"
    "\"What can be said at all can be said clearly; and of what one cannot talk, about that one must be silent.\"\n"
    "- Ludwig Wittgenstein\n\n"
    "Key Points:\n"
    "1. Emphasizes the importance of clear expression in language.\n"
    "2. Suggests that unclear or indescribable concepts should not be discussed.\n"
    "3. Reflects on the limits of language and thought."
)

# Slide 4: Combining the Perspectives
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.placeholders[1]
title.text = "Combining the Perspectives"
content.text = (
    "1. Language and Computation:\n"
    "   - Clear language is essential for defining computational problems.\n"
    "   - Only clearly defined problems can be addressed by machines.\n\n"
    "2. Expression and Problem-Solving:\n"
    "   - von Neumann's confidence relies on clear problem definition.\n"
    "   - Wittgenstein's clarity principle supports the foundation of computational theory.\n\n"
    "3. Philosophical and Practical Implications:\n"
    "   - Theoretical limits (e.g., the Halting Problem) vs. practical solutions.\n"
    "   - Continuous technological progress expands the boundaries of what machines can do."
)

# Slide 5: Conclusion
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]
title.text = "Conclusion"
content.text = (
    "1. Wittgenstein and von Neumann provide complementary views on language and computation.\n"
    "2. Clear expression is crucial for both philosophical clarity and computational feasibility.\n"
    "3. The interplay between theoretical limits and practical capabilities drives technological advancement."
)

# Save the presentation
pptx_path = "/mnt/data/VonNeumann_Wittgenstein_Perspectives_New.pptx"
prs.save(pptx_path)
